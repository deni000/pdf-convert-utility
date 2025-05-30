import os
import platform
from PyPDF2 import PdfReader, PdfWriter, PdfMerger
import pdfplumber
from pdf2docx import Converter
from pptx import Presentation
from pptx.util import Inches

# Function definitions
def pdf_to_word(pdf_path, docx_path):
    cv = Converter(pdf_path)
    cv.convert(docx_path)
    cv.close()
    print("PDF converted to Word successfully. Result saved as", docx_path)

def pdf_to_powerpoint(pdf_path, pptx_path):
    prs = Presentation()
    with pdfplumber.open(pdf_path) as pdf:
        for page in pdf.pages:
            image_path = "temp_image.png"
            page.to_image(resolution=300).save(image_path)
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.add_picture(image_path, 0, 0, width=Inches(10), height=Inches(7.5))
            os.remove(image_path)
    prs.save(pptx_path)
    print("PDF converted to PowerPoint successfully. Result saved as", pptx_path)

def merge_pdf(pdf_files, output_path):
    merger = PdfMerger()
    for pdf_file in pdf_files:
        merger.append(pdf_file)
    with open(output_path, "wb") as f_out:
        merger.write(f_out)
    merger.close()
    print("PDFs merged successfully. Result saved as", output_path)

def shrink_pdf(input_pdf, output_pdf):
    with pdfplumber.open(input_pdf) as pdf:
        pdf.pages[0].to_image()  # Forces image decompression

    # Ensure output path ends with .pdf
    if not output_pdf.lower().endswith('.pdf'):
        output_pdf += '.pdf'

    # Use Ghostscript to compress PDF, suppressing warnings
    if platform.system() == "Windows":
        cmd = (
            f"gswin64c -sDEVICE=pdfwrite -dCompatibilityLevel=1.4 -dPDFSETTINGS=/ebook "
            f"-dNOPAUSE -dQUIET -dBATCH -sOutputFile=\"{output_pdf}\" \"{input_pdf}\" > nul 2>&1"
        )
    else:
        cmd = (
            f"gs -sDEVICE=pdfwrite -dCompatibilityLevel=1.4 -dPDFSETTINGS=/ebook "
            f"-dNOPAUSE -dQUIET -dBATCH -sOutputFile='{output_pdf}' '{input_pdf}' 2>/dev/null"
        )

    os.system(cmd)
    print("PDF compressed successfully. Result saved as", output_pdf)

def split_pdf(input_pdf, output_dir):
    pdf_reader = PdfReader(input_pdf)
    for page_number, page in enumerate(pdf_reader.pages):
        pdf_writer = PdfWriter()
        pdf_writer.add_page(page)
        output_pdf_path = os.path.join(output_dir, f"page_{page_number + 1}.pdf")
        with open(output_pdf_path, "wb") as output_pdf:
            pdf_writer.write(output_pdf)
    print("PDF split successfully. Result saved in", output_dir)

# Main function
def main():
    print("Available functions:")
    print("1. PDF to Word")
    print("2. PDF to PowerPoint")
    print("3. Merge PDF")
    print("4. Compress PDF")
    print("5. Split PDF")
    choice = input("Enter the number corresponding to the function you want to use: ")

    if choice == "1":
        pdf_path = input("Enter the path of the PDF file to convert to Word: ").strip()
        docx_path = input("Enter the path for the output Word file (leave blank for default): ").strip()
        if not docx_path:
            docx_path = pdf_path.replace('.pdf', '.docx')
        pdf_to_word(pdf_path, docx_path)

    elif choice == "2":
        pdf_path = input("Enter the path of the PDF file to convert to PowerPoint: ").strip()
        pptx_path = input("Enter the path for the output PowerPoint file (leave blank for default): ").strip()
        if not pptx_path:
            pptx_path = pdf_path.replace('.pdf', '.pptx')
        pdf_to_powerpoint(pdf_path, pptx_path)

    elif choice == "3":
        pdf_files_to_merge = []
        while True:
            pdf_path = input("Enter the path of the PDF file to merge (or enter 'done' to finish): ").strip()
            if pdf_path.lower() == 'done':
                break
            pdf_files_to_merge.append(pdf_path)
        output_path = input("Enter the output path for the merged PDF: ").strip()
        if not output_path.lower().endswith('.pdf'):
            output_path += '.pdf'
        merge_pdf(pdf_files_to_merge, output_path)

    elif choice == "4":
        input_pdf_to_compress = input("Enter the path of the input PDF file to compress: ").strip()
        output_compressed_pdf = input("Enter the output path for the compressed PDF: ").strip()
        shrink_pdf(input_pdf_to_compress, output_compressed_pdf)

    elif choice == "5":
        input_pdf_to_split = input("Enter the path of the input PDF file to split: ").strip()
        output_dir_for_split_pages = input("Enter the output directory for split PDF pages: ").strip()
        os.makedirs(output_dir_for_split_pages, exist_ok=True)
        split_pdf(input_pdf_to_split, output_dir_for_split_pages)

    else:
        print("Invalid choice. Please enter a valid number.")

if __name__ == "__main__":
    main()
